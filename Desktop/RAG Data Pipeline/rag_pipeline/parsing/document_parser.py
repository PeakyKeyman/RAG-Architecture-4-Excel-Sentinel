"""
Document parser for executive document types.
Handles PDFs, DOCX, PPTX with text, tables, and metadata extraction.
"""

import os
import io
from typing import Dict, Any, List, Optional, Union
from pathlib import Path
from dataclasses import dataclass
import logging

# Document parsing libraries
try:
    import PyPDF2
    import pdfplumber
    from docx import Document as DocxDocument
    from pptx import Presentation
    import pandas as pd
    from google.cloud import storage
except ImportError as e:
    logging.warning(f"Document parsing dependencies not installed: {e}")

from ..models.chunk import Document
from ..core.logging_config import get_logger


logger = get_logger(__name__, "document_parser")


@dataclass
class ParsedContent:
    """Container for parsed document content."""
    text: str
    tables: List[Dict[str, Any]]
    metadata: Dict[str, Any]
    images: List[Dict[str, Any]]  # For future image processing
    

class DocumentTypeError(Exception):
    """Raised when document type is not supported."""
    pass


class DocumentParser:
    """
    Parses executive document types with multi-tenant security.
    
    Supported formats:
    - PDF: Financial reports, board presentations, strategic documents
    - DOCX: Policy documents, communications, frameworks
    - PPTX: Investor presentations, strategy decks
    """
    
    SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.pptx'}
    
    def __init__(self, 
                 storage_client: Optional[storage.Client] = None,
                 temp_dir: str = "/tmp/rag_parser"):
        """
        Initialize document parser.
        
        Args:
            storage_client: Google Cloud Storage client
            temp_dir: Directory for temporary file processing
        """
        self.storage_client = storage_client or storage.Client()
        self.temp_dir = Path(temp_dir)
        self.temp_dir.mkdir(exist_ok=True)
        
    def parse_document(self, 
                      file_path: Union[str, Path], 
                      user_id: str,
                      group_id: str, 
                      org_id: str,
                      document_metadata: Optional[Dict[str, Any]] = None) -> Document:
        """
        Parse document with multi-tenant metadata.
        
        Args:
            file_path: Path to document (local or GCS URI)
            user_id: User identifier for access control
            group_id: Group identifier for access control  
            org_id: Organization identifier for access control
            document_metadata: Additional document metadata
            
        Returns:
            Document object with parsed content and security metadata
        """
        try:
            file_path = Path(file_path) if not str(file_path).startswith('gs://') else file_path
            
            # Validate file type
            if isinstance(file_path, Path):
                extension = file_path.suffix.lower()
            else:
                extension = Path(str(file_path)).suffix.lower()
                
            if extension not in self.SUPPORTED_EXTENSIONS:
                raise DocumentTypeError(f"Unsupported file type: {extension}")
            
            # Download from GCS if needed
            local_path = self._ensure_local_file(file_path)
            
            # Parse based on file type
            if extension == '.pdf':
                parsed_content = self._parse_pdf(local_path)
            elif extension == '.docx':
                parsed_content = self._parse_docx(local_path)
            elif extension == '.pptx':
                parsed_content = self._parse_pptx(local_path)
            
            # Create document with security metadata
            document = self._create_document(
                content=parsed_content,
                file_path=str(file_path),
                user_id=user_id,
                group_id=group_id,
                org_id=org_id,
                document_metadata=document_metadata or {}
            )
            
            # Cleanup temp file if downloaded
            if str(file_path).startswith('gs://') and local_path.exists():
                local_path.unlink()
                
            logger.info(f"Successfully parsed document: {file_path}", 
                       extra={"org_id": org_id, "user_id": user_id})
            
            return document
            
        except Exception as e:
            logger.error(f"Failed to parse document {file_path}: {str(e)}", 
                        exc_info=True, extra={"org_id": org_id, "user_id": user_id})
            raise
    
    def _ensure_local_file(self, file_path: Union[str, Path]) -> Path:
        """Download file from GCS if needed, return local path."""
        if str(file_path).startswith('gs://'):
            # Parse GCS URI
            gcs_path = str(file_path)[5:]  # Remove 'gs://'
            bucket_name, blob_name = gcs_path.split('/', 1)
            
            # Download to temp directory
            local_filename = Path(blob_name).name
            local_path = self.temp_dir / local_filename
            
            bucket = self.storage_client.bucket(bucket_name)
            blob = bucket.blob(blob_name)
            blob.download_to_filename(local_path)
            
            return local_path
        else:
            return Path(file_path)
    
    def _parse_pdf(self, file_path: Path) -> ParsedContent:
        """Parse PDF document extracting text and tables."""
        text_content = []
        tables = []
        metadata = {}
        
        # Use pdfplumber for better table extraction
        with pdfplumber.open(file_path) as pdf:
            metadata['page_count'] = len(pdf.pages)
            metadata['pdf_metadata'] = pdf.metadata
            
            for page_num, page in enumerate(pdf.pages, 1):
                # Extract text
                page_text = page.extract_text()
                if page_text:
                    text_content.append(f"[Page {page_num}]\n{page_text}")
                
                # Extract tables
                page_tables = page.extract_tables()
                for table_num, table in enumerate(page_tables):
                    if table and len(table) > 1:  # Skip empty or single-row tables
                        tables.append({
                            'page': page_num,
                            'table_number': table_num + 1,
                            'headers': table[0] if table[0] else [],
                            'rows': table[1:],
                            'markdown': self._table_to_markdown(table)
                        })
        
        return ParsedContent(
            text='\n\n'.join(text_content),
            tables=tables,
            metadata=metadata,
            images=[]  # PDF image extraction for future enhancement
        )
    
    def _parse_docx(self, file_path: Path) -> ParsedContent:
        """Parse DOCX document extracting text and tables."""
        doc = DocxDocument(file_path)
        
        text_content = []
        tables = []
        metadata = {
            'paragraph_count': len(doc.paragraphs),
            'table_count': len(doc.tables)
        }
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text)
        
        # Extract tables
        for table_num, table in enumerate(doc.tables, 1):
            table_data = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                table_data.append(row_data)
            
            if table_data and len(table_data) > 1:
                tables.append({
                    'table_number': table_num,
                    'headers': table_data[0] if table_data[0] else [],
                    'rows': table_data[1:],
                    'markdown': self._table_to_markdown(table_data)
                })
        
        return ParsedContent(
            text='\n\n'.join(text_content),
            tables=tables,
            metadata=metadata,
            images=[]
        )
    
    def _parse_pptx(self, file_path: Path) -> ParsedContent:
        """Parse PPTX presentation extracting text from slides."""
        prs = Presentation(file_path)
        
        text_content = []
        tables = []
        metadata = {
            'slide_count': len(prs.slides),
            'presentation_title': 'Unknown'
        }
        
        # Try to get presentation title from first slide
        if prs.slides:
            first_slide = prs.slides[0]
            for shape in first_slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    metadata['presentation_title'] = shape.text.strip()
                    break
        
        # Extract text from all slides
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Extract tables from slides
                if shape.has_table:
                    table_data = []
                    table = shape.table
                    for row in table.rows:
                        row_data = [cell.text.strip() for cell in row.cells]
                        table_data.append(row_data)
                    
                    if table_data:
                        tables.append({
                            'slide': slide_num,
                            'headers': table_data[0] if len(table_data) > 1 else [],
                            'rows': table_data[1:] if len(table_data) > 1 else table_data,
                            'markdown': self._table_to_markdown(table_data)
                        })
            
            if len(slide_text) > 1:  # More than just the slide header
                text_content.append('\n'.join(slide_text))
        
        return ParsedContent(
            text='\n\n'.join(text_content),
            tables=tables,
            metadata=metadata,
            images=[]
        )
    
    def _table_to_markdown(self, table_data: List[List[str]]) -> str:
        """Convert table data to markdown format."""
        if not table_data or not table_data[0]:
            return ""
        
        markdown = []
        
        # Headers
        headers = table_data[0]
        markdown.append('| ' + ' | '.join(headers) + ' |')
        markdown.append('| ' + ' | '.join(['---'] * len(headers)) + ' |')
        
        # Rows
        for row in table_data[1:]:
            # Pad row to match header length
            padded_row = row + [''] * (len(headers) - len(row))
            markdown.append('| ' + ' | '.join(padded_row[:len(headers)]) + ' |')
        
        return '\n'.join(markdown)
    
    def _create_document(self, 
                        content: ParsedContent,
                        file_path: str,
                        user_id: str,
                        group_id: str,
                        org_id: str,
                        document_metadata: Dict[str, Any]) -> Document:
        """Create Document object with parsed content and security metadata."""
        
        # Generate document ID
        document_id = f"{org_id}_{Path(file_path).stem}_{hash(file_path) % 100000}"
        
        # Combine text and tables
        full_content = [content.text]
        
        if content.tables:
            full_content.append("\n\n--- EXTRACTED TABLES ---\n")
            for i, table in enumerate(content.tables, 1):
                full_content.append(f"\nTable {i}:\n{table['markdown']}")
        
        # Security and document metadata
        metadata = {
            # Security metadata (critical for multi-tenant filtering)
            'user_id': user_id,
            'group_id': group_id, 
            'org_id': org_id,
            'access_level': document_metadata.get('access_level', 'standard'),
            
            # Document metadata
            'source_path': file_path,
            'file_type': Path(file_path).suffix.lower(),
            'document_category': document_metadata.get('category', 'general'),
            'document_type': self._classify_document_type(file_path, content),
            'processed_timestamp': pd.Timestamp.now().isoformat(),
            
            # Content metadata
            'has_tables': len(content.tables) > 0,
            'table_count': len(content.tables),
            **content.metadata,
            **document_metadata
        }
        
        return Document(
            document_id=document_id,
            content='\n\n'.join(full_content),
            metadata=metadata,
            source_path=file_path
        )
    
    def _classify_document_type(self, file_path: str, content: ParsedContent) -> str:
        """Classify document type based on filename and content."""
        filename = Path(file_path).name.lower()
        text_lower = content.text.lower()
        
        # Financial documents
        if any(term in filename for term in ['financial', 'p&l', 'balance', 'income', 'cash_flow']):
            return 'financial'
        if any(term in text_lower for term in ['revenue', 'ebitda', 'gross margin', 'operating income']):
            return 'financial'
            
        # Strategic documents
        if any(term in filename for term in ['strategy', 'strategic', 'vision', 'mission', 'okr']):
            return 'strategic'
        if any(term in text_lower for term in ['strategic initiative', 'vision', 'mission', 'objectives']):
            return 'strategic'
            
        # Board/Investor documents
        if any(term in filename for term in ['board', 'investor', 'presentation', 'deck']):
            return 'board_investor'
        if 'board of directors' in text_lower or 'investor' in text_lower:
            return 'board_investor'
            
        # Policy documents
        if any(term in filename for term in ['policy', 'compliance', 'governance', 'risk']):
            return 'policy_compliance'
            
        # Market research
        if any(term in filename for term in ['market', 'research', 'analysis', 'competitive']):
            return 'market_research'
            
        return 'general'